from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def create_section_header(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    return slide

def create_content_slide(prs, title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def create_two_column_slide(prs, title, left_items, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    left = 0.5
    top = 1.5
    width = 4.5
    height = 5
    
    left_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(6)
    
    right_box = slide.shapes.add_textbox(Inches(left + width + 0.2), Inches(top), Inches(width), Inches(height))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(6)
    
    return slide

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

create_title_slide(prs, 
    "UAV-based Traffic & Pedestrian Monitoring",
    "Onderzoeksvragen Analyse\n569 Papers (2001-2026)\n\nDavy Janssens\nMaart 2026"
)

create_content_slide(prs, "📊 Dataset Overzicht", [
    "• Totaal aantal papers: 569",
    "• Tijdsperiode: 2001-2026",
    "• Hoofdthema's:",
    "  - UAV-based detection & tracking",
    "  - Parking monitoring",
    "  - Pedestrian analysis",
    "  - Traffic flow & congestion",
    "  - Deep learning methodologieën"
])

create_section_header(prs, "1️⃣ METHODOLOGISCH\nDeep Learning & Algoritmes")

create_content_slide(prs, "1.1 Object Detection Architecturen", [
    "Centrale vraag:",
    "Welke deep learning architecturen zijn het meest geschikt voor UAV-based object detection?",
    "",
    "Key bevindingen:",
    "• 35+ papers over YOLO-implementaties",
    "• YOLOv8 en YOLOv10 met knowledge distillation: beste real-time prestaties",
    "• SSD: sneller maar minder accuraat bij kleine objecten",
    "• Lightweight models (MobileNet): geschikt voor edge devices"
])

create_content_slide(prs, "1.2 Multi-Scale Feature Fusion", [
    "Centrale vraag:",
    "Hoe kunnen multi-scale features optimaal worden gefuseerd voor kleine object detectie?",
    "",
    "Deelvragen:",
    "• Welke feature pyramid architecturen werken het best?",
    "• Hoe kunnen attention mechanisms detectie verbeteren?",
    "  - Spatial attention",
    "  - Channel attention",
    "• Optimale receptive field sizes voor verschillende vlieghoogtes?"
])

create_content_slide(prs, "1.3 Tracking Algoritmes", [
    "Centrale vraag:",
    "Welke tracking methoden zijn robuust genoeg voor bewegende camera platforms?",
    "",
    "Vergelijking:",
    "• SORT: Basis tracking",
    "• DeepSORT: Met appearance features",
    "• ByteTrack: State-of-the-art voor UAV",
    "",
    "Uitdagingen:",
    "• Camera motion compensation",
    "• Occlusion handling",
    "• Trajectory-based methods"
])

create_section_header(prs, "2️⃣ TECHNISCH\nReal-time Processing & Hardware")

create_content_slide(prs, "2.1 Real-time Verwerking", [
    "Centrale vraag:",
    "Wat zijn de bottlenecks voor real-time UAV video processing?",
    "",
    "Bevindingen uit papers:",
    "• Minimale frame rate:",
    "  - 1/3 Hz voor visuele interpretatie",
    "  - 3+ Hz voor automatische tracking",
    "",
    "• NVIDIA Jetson TX2: 4K real-time processing mogelijk",
    "• CPU + GPU parallel: latency verbergen"
])

create_content_slide(prs, "2.2 Edge Computing", [
    "Centrale vraag:",
    "Kunnen UAVs autonome on-board processing uitvoeren?",
    "",
    "Embedded platforms:",
    "• NVIDIA Jetson (TX2, Nano, Orin)",
    "• Google Coral TPU",
    "• Intel Neural Compute Stick",
    "",
    "Optimalisatie technieken:",
    "• Model compression (pruning, quantization)",
    "• Battery constraints vs computational capabilities"
])

create_content_slide(prs, "2.3 Image Registration & Stabilization", [
    "Centrale vraag:",
    "Hoe kan camera motion worden gecompenseerd voor accurate metingen?",
    "",
    "Technieken:",
    "• SIFT: Scale-Invariant Feature Transform",
    "• SURF: Speeded-Up Robust Features",
    "• Optical flow",
    "• Homography estimation",
    "",
    "Toepassing:",
    "• Speed estimation zonder constante drone snelheid",
    "• World coordinate transformation"
])

create_section_header(prs, "3️⃣ TOEPASSING\nUse Cases")

create_content_slide(prs, "3.1 Parking Management", [
    "Centrale vraag:",
    "Hoe kunnen UAVs parking occupancy monitoring revolutioneren?",
    "",
    "Bevindingen:",
    "• Optimale vlieghoogte: ~50m",
    "  - Balans tussen coverage en resolutie",
    "",
    "• License plate recognition mogelijk (ALPR)",
    "• Route optimization via TSP voor meerdere parking lots",
    "",
    "Papers: 10 papers over parking detection"
])

create_content_slide(prs, "3.2 Pedestrian Behavior Analysis", [
    "Centrale vraag:",
    "Welke inzichten kunnen UAVs bieden over voetgangersgedrag?",
    "",
    "Toepassingen:",
    "• Walking speed metingen op street-scale",
    "  - YOLOv5 + DeepSORT + geometric correction",
    "",
    "• Pedestrian flow patterns op shopping streets",
    "• Pedestrian-vehicle conflict detection",
    "  - Surrogate safety measures (TTC, PET)",
    "",
    "Papers: 124 papers over pedestrian analysis"
])

create_content_slide(prs, "3.3 Traffic Flow & Congestion", [
    "Centrale vraag:",
    "Hoe effectief zijn UAVs voor traffic monitoring vs traditionele methoden?",
    "",
    "Meetbare parameters:",
    "• Volume (vehicle counting)",
    "• Speed (snelheidsmetingen)",
    "• Density (voertuigdichtheid)",
    "",
    "Geavanceerde analyses:",
    "• Origin-destination matrices",
    "• Traffic violation detection",
    "  - Illegal lane changes, double parking, crosswalk obstructions"
])

create_content_slide(prs, "3.4 Roundabout & Intersection Analysis", [
    "Centrale vraag:",
    "Wat is de meerwaarde van UAVs voor complexe verkeerssituaties?",
    "",
    "openDD Dataset (Breuer et al., 2020):",
    "• 84,774 nauwkeurig getrackte trajectories",
    "• 7 verschillende rotondes",
    "• 62 uur aan trajectory data",
    "",
    "Analyses:",
    "• Traffic conflict identification",
    "• Interaction patterns tussen weggebruikers",
    "• Yielding behavior"
])

create_section_header(prs, "4️⃣ ECONOMISCH & PRAKTISCH")

create_content_slide(prs, "4.1 Cost-Benefit Analysis", [
    "Centrale vraag:",
    "Wat is de economische haalbaarheid van UAV-based monitoring?",
    "",
    "Bevindingen (Jairus et al., 2025):",
    "• 60-80% kostenreductie vs traditionele methoden",
    "• 8-point density threshold voor viability",
    "",
    "Impact operational intensity:",
    "• 30 dagen/jaar: 9.0% viability",
    "• 180 dagen/jaar: 81.6% viability",
    "",
    "Conclusie: Economisch voordelig voor 67% van urban areas"
])

create_content_slide(prs, "4.2 Deployment & Scalability", [
    "Centrale vraag:",
    "Hoe kunnen UAV-systemen op grote schaal worden ingezet?",
    "",
    "Uitdagingen:",
    "• Regelgeving in urban areas",
    "• Privacy implications van aerial surveillance",
    "",
    "Enablers:",
    "• 5G/nomadic networks voor UAV operations",
    "• Cellular connectivity voor real-time data streaming",
    "• Cloud processing voor heavy computations"
])

create_content_slide(prs, "4.3 Data Quality & Accuracy", [
    "Centrale vraag:",
    "Hoe accuraat zijn UAV-metingen vergeleken met ground truth?",
    "",
    "Detection accuracy:",
    "• 95-100% detection rates haalbaar",
    "  - Vehicles, cyclists, pedestrians",
    "",
    "Variatie factoren:",
    "• Weersomstandigheden (day/night, rain, fog)",
    "• Altitude en viewing angle",
    "• Object size en occlusion"
])

create_section_header(prs, "5️⃣ DATA & DATASETS")

create_content_slide(prs, "5.1 Beschikbare Datasets", [
    "Belangrijkste datasets:",
    "",
    "• VisDrone2019: Urban traffic, vehicles & pedestrians",
    "• UAVDT: UAV detection and tracking",
    "• CARPK: 90,000 cars in parking lots",
    "• openDD: 84,774 trajectories in roundabouts",
    "• UA-DETRAC: Urban road environments",
    "",
    "Papers: 32 papers over datasets"
])

create_content_slide(prs, "5.2 Dataset Gaps & Synthetic Data", [
    "Ondervertegenwoordigd:",
    "• Extreme weather conditions",
    "• Night-time scenarios",
    "• Multi-modal interactions",
    "• Long-term temporal patterns",
    "",
    "Synthetic datasets:",
    "• AirSim (Microsoft)",
    "• Unreal Engine simulations",
    "• Data augmentation strategies"
])

create_section_header(prs, "6️⃣ MULTI-MODAL & SENSOR FUSION")

create_content_slide(prs, "6.1 Thermal vs RGB Imaging", [
    "Centrale vraag:",
    "Wanneer is thermal imaging superieur aan RGB?",
    "",
    "Thermal voordelen:",
    "• Night-time pedestrian detection",
    "• Prominent thermal signature van mensen",
    "• Werkt bij slechte verlichting",
    "",
    "RGB + Thermal fusion:",
    "• Betere accuracy door complementaire informatie",
    "• Trade-off: sensor cost vs performance"
])

create_content_slide(prs, "6.2 Multi-Sensor Integration", [
    "Centrale vraag:",
    "Hoe kunnen verschillende sensoren worden gecombineerd?",
    "",
    "Sensor combinaties:",
    "• LiDAR + Camera: 3D trajectory reconstruction",
    "• GPS/IMU: Image registration ondersteuning",
    "• Hyperspectral: Geavanceerde traffic analysis",
    "",
    "Voordelen:",
    "• Robustere detectie",
    "• Betere accuracy",
    "• Redundantie bij sensor failure"
])

create_section_header(prs, "7️⃣ EMERGING TRENDS")

create_content_slide(prs, "7.1 Swarm Intelligence", [
    "Centrale vraag:",
    "Kunnen drone swarms collaborative monitoring uitvoeren?",
    "",
    "Mogelijkheden:",
    "• Coverage optimalisatie door meerdere UAVs",
    "• Occlusion oplossen via multi-view fusion",
    "• Distributed processing",
    "",
    "Uitdagingen:",
    "• Communicatie requirements",
    "• Swarm coordination algoritmes",
    "• Collision avoidance"
])

create_content_slide(prs, "7.2 Predictive Analytics", [
    "Centrale vraag:",
    "Kunnen UAV-data worden gebruikt voor traffic prediction?",
    "",
    "Technieken:",
    "• LSTM/Transformer models voor trajectory forecasting",
    "• Congestion pattern prediction uit historical data",
    "• Digital twin development voor smart cities",
    "",
    "Toepassingen:",
    "• Proactive traffic management",
    "• Incident prediction",
    "• Route optimization"
])

create_content_slide(prs, "7.3 Autonomous Operations", [
    "Centrale vraag:",
    "Hoe autonoom kunnen UAV monitoring systemen worden?",
    "",
    "Autonomie levels:",
    "• Event detection: Automatisch interessante events detecteren",
    "• Adaptive focus: Zelfstandig op events focussen",
    "• Path planning: Reinforcement learning voor optimal routes",
    "",
    "Safety implications:",
    "• Fully autonomous operations in urban areas",
    "• Fail-safe mechanisms",
    "• Human-in-the-loop vs full autonomy"
])

create_section_header(prs, "🔑 RESEARCH GAPS")

create_two_column_slide(prs, "Belangrijkste Gaps in Huidig Onderzoek",
    [
        "METHODOLOGISCH:",
        "• Cross-domain generalization",
        "• Extreme weather performance",
        "• Nighttime detection",
        "",
        "TECHNISCH:",
        "• Battery life vs computational load",
        "• 5G/6G integration",
        "• Privacy-preserving techniques"
    ],
    [
        "TOEPASSING:",
        "• Long-term deployment studies",
        "• Multi-modal traffic analysis",
        "• Incident detection & response",
        "",
        "ECONOMISCH:",
        "• Total Cost of Ownership",
        "• Regulatory compliance costs",
        "• Maintenance & operational costs"
    ]
)

create_section_header(prs, "💡 MEEST VEELBELOVENDE RICHTINGEN")

create_content_slide(prs, "Top 5 Research Opportunities", [
    "1. Unified Multi-Task Framework",
    "   → End-to-end learning voor alle taken simultaan",
    "",
    "2. Edge AI Optimization",
    "   → Hardware-software co-design voor optimal performance",
    "",
    "3. Economic Viability Studies",
    "   → Large-scale cost-benefit analyses",
    "",
    "4. Privacy-Preserving Monitoring",
    "   → On-device anonymization, federated learning",
    "",
    "5. Autonomous Adaptive Monitoring",
    "   → Event-driven intelligent monitoring"
])

create_section_header(prs, "📈 TIJDLIJN VAN ONTWIKKELINGEN")

create_content_slide(prs, "Evolutie 2001-2026", [
    "2001-2010:",
    "• Fundamentele traffic dynamics",
    "• Eerste aerial monitoring experimenten",
    "",
    "2016-2018:",
    "• Opkomst deep learning (YOLO, Faster R-CNN)",
    "• Eerste UAV-specific detection papers",
    "",
    "2019-2021:",
    "• Focus op real-time processing",
    "• Edge computing toepassingen",
    "",
    "2022-2024:",
    "• Multi-task frameworks",
    "• Attention mechanisms",
    "",
    "2025-2026:",
    "• Knowledge distillation",
    "• Economic feasibility studies",
    "• Unified intelligent systems"
])

create_section_header(prs, "📚 TOP PAPERS")

create_content_slide(prs, "Must-Read Papers (1/2)", [
    "1. Khan et al. (2026) - Unified Framework for Vehicle Detection",
    "   → YOLOv10 + Knowledge Distillation + Cross-domain",
    "",
    "2. Khanpour et al. (2025) - UAV-Based Intelligent Surveillance",
    "   → Real-time detection, tracking, behavioral analysis",
    "",
    "3. Alnusayri et al. (2026) - Traffic Vision",
    "   → LSTM trajectory forecasting + pattern analysis",
    "",
    "4. Breuer et al. (2020) - openDD Dataset",
    "   → 84,774 trajectories, interaction patterns"
])

create_content_slide(prs, "Must-Read Papers (2/2)", [
    "5. Jairus et al. (2025) - Economic Feasibility",
    "   → 60-80% cost reduction, viability analysis",
    "",
    "6. Peraza-Garzón et al. (2026) - YOLOv11 Parking",
    "   → Latest YOLO for parking occupancy",
    "",
    "7. Jiao & Fei (2023) - Pedestrian Walking Speed",
    "   → Street-scale behavior extraction",
    "",
    "8. Gettman & Head (2003) - Surrogate Safety Measures",
    "   → Framework voor conflict detection (TTC, PET)"
])

create_section_header(prs, "🎯 CONCLUSIES")

create_content_slide(prs, "Belangrijkste Bevindingen", [
    "METHODOLOGISCH:",
    "• YOLOv8/v10 + knowledge distillation: state-of-the-art",
    "• Multi-scale feature fusion essentieel voor kleine objecten",
    "",
    "TECHNISCH:",
    "• Real-time processing mogelijk (Jetson TX2, 4K @ 30fps)",
    "• 3+ Hz minimaal voor automatische tracking",
    "",
    "TOEPASSING:",
    "• 95-100% detection rates haalbaar",
    "• Optimale vlieghoogte: ~50m",
    "",
    "ECONOMISCH:",
    "• 60-80% kostenreductie mogelijk",
    "• Viability bij 180 dagen/jaar: 81.6%"
])

create_content_slide(prs, "Aanbevelingen voor Toekomstig Onderzoek", [
    "METHODOLOGISCH:",
    "• Benchmark datasets voor extreme conditions",
    "• Domain adaptation technieken",
    "• Self-supervised learning",
    "",
    "TECHNISCH:",
    "• Hardware-algorithm co-optimization",
    "• Energy-efficient inference strategies",
    "• 5G/6G integration",
    "",
    "TOEPASSING:",
    "• Longitudinal studies (>6 maanden)",
    "• Multi-city comparative analyses",
    "• Integration met existing systems"
])

create_content_slide(prs, "Bronnen & Referenties", [
    "Analyse gebaseerd op:",
    "• 569 papers uit Bilal Drones Data References",
    "• Tijdsperiode: 2001-2026",
    "• Hoofdcategorieën:",
    "  - Detection & Tracking (78 papers)",
    "  - Pedestrian Analysis (124 papers)",
    "  - Traffic Flow (46 papers)",
    "  - Speed Estimation (69 papers)",
    "  - Deep Learning (35 papers)",
    "  - Real-time Processing (30 papers)",
    "  - Datasets (32 papers)",
    "  - Parking Management (10 papers)"
])

create_title_slide(prs,
    "Vragen?",
    "Davy Janssens\ndavy.janssens@uhasselt.be\n\nGitHub: github.com/davyjanssens-dotcom/dronespapers"
)

prs.save('UAV_Traffic_Monitoring_Presentatie.pptx')
print("✅ Presentatie succesvol aangemaakt: UAV_Traffic_Monitoring_Presentatie.pptx")
print(f"📊 Totaal aantal slides: {len(prs.slides)}")
